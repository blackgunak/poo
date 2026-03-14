"""
Génère un fichier PowerPoint académique décrivant l'application Cinema.
Usage:
  pip install python-pptx
  python presentation_generate.py
Produira `Cinema_Presentation.pptx` dans le dossier courant.

Le script insère des slides : titre, plan, contexte, architecture, modèles, fonctionnalités, flux QR, migrations, admin, UX, démo (placeholders), conclusion.
Remplace/ajoute des images dans `media/` pour les inclure automatiquement si trouvées.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path

OUTFILE = 'Cinema_Presentation.pptx'
prs = Presentation()

# Helper styles
def set_title(slide, title, subtitle=None):
    title_box = slide.shapes.title
    title_box.text = title
    title_tf = title_box.text_frame
    title_tf.paragraphs[0].font.size = Pt(36)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(20,20,20)
    if subtitle:
        # add subtitle in a content placeholder if available
        for shp in slide.placeholders:
            if shp.placeholder_format.type.name == 'BODY':
                shp.text = subtitle
                shp.text_frame.paragraphs[0].font.size = Pt(14)
                break


def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_title(slide, 'Cinema — Application de gestion de réservation', 'Présentation académique')
slide.placeholders[1].text = 'Auteur: Équipe projet\nDate: 2026-03-05\nContexte: Projet Django — Système de réservations et QR single-use'

# Outline
add_bullet_slide('Plan', [
    'Contexte & objectifs',
    'Architecture technique',
    'Modèles de données clés',
    'Fonctionnalités principales',
    'Flux QR codes single-use',
    'Migrations et contraintes',
    'Interface admin & staff',
    'Démonstration et captures',
    'Conclusion & perspectives'
])

# Contexte
add_bullet_slide('Contexte & Objectifs', [
    "Application Django pour gestion d'un cinéma (films, séances, réservations)",
    "Objectifs: UX simple, QR single-use pour validation d'entrée, interface staff/admin",
    "Base de données: SQLite (dev), migration prudente pour champs uniques"
])

# Architecture
add_bullet_slide('Architecture technique', [
    'Backend: Django (v6+)',
    'Frontend: Templates Django + CSS/JS (seatmap, modals)',
    'QR generation: package Python `qrcode` -> images PNG',
    'Auth: Django auth, logout POST-only, staff vs user roles'
])

# Models
add_bullet_slide('Modèles de données', [
    'Film, Salle, Séance, Reservation',
    'Reservation: champs QR -> `qr_token` (UUID), `qr_valid`, `qr_used_at`, `qr_used_by`',
    'Stratégie: token unique, invalidation single-use, champ timestamp + FK utilisateur'
])

# Features
add_bullet_slide('Fonctionnalités principales', [
    'Listing films, page d’accueil avec hero “À l\'affiche maintenant”',
    'Sélection de sièges via seatmap JS',
    'Génération & téléchargement QR pour chaque réservation',
    'Scan public (single-use) + endpoints staff/admin (AJAX)',
    'Migrations sûres pour ajout de UUID uniques sur SQLite'
])

# QR flow
add_bullet_slide('Flux QR single-use', [
    'Réservation -> génération `qr_token` (UUID) + image PNG',
    'QRCode rendu propriétaire téléchargeable; visualisable en modal',
    'Endpoint public de scan: lecture du token -> si valide, marquer utilisé et retourner résultat',
    'Staff/admin disposent d\'endpoints AJAX pour scan sans redirection'
])

# Migrations
add_bullet_slide('Migrations & sécurité', [
    'Problème SQLite: ajout direct d\'une colonne unique cause IntegrityError',
    'Solution: add nullable field -> populate UUIDs -> alter vers non-null + unique',
    'Pratique: scripts RunPython pour backfill en migration'
])

# Admin & staff
add_bullet_slide('Admin et interface staff', [
    'ReservationAdmin: vue de scan, change_list hook',
    'Custom staff reservations list sur le site pour scan AJAX',
    'Permissions: owner-only endpoints et accès staff pour actions spécifiques'
])

# Demo / captures (placeholders)
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_title(slide, 'Démonstration (captures d\'écran)')
shapes = slide.shapes
img_dir = Path('media/affiches')
# Try to insert up to 2 images if present
imgs = list(img_dir.glob('*.*'))[:2] if img_dir.exists() else []
if imgs:
    left = Inches(0.5)
    top = Inches(1.6)
    for img in imgs:
        try:
            shapes.add_picture(str(img), left, top, width=Inches(4.2))
            left += Inches(4.6)
        except Exception:
            pass
else:
    # add bullet placeholders
    tf = shapes.placeholders[1].text_frame
    tf.text = 'Insérer captures: page d\'accueil, réservation, modal QR, interface staff.'
    tf.paragraphs[0].font.size = Pt(16)

# UX et design
add_bullet_slide('UX & Design', [
    'Hero large sur la page d\'accueil: mise en avant du film du moment',
    'Cartes film plus lisibles (affiches agrandies)',
    'Toasts/modals pour retours utilisateurs, style SaaS moderne'
])

# Limitations & perspectives
add_bullet_slide('Limites & Perspectives', [
    'Passage SQLite -> PostgreSQL en production',
    'Auditing + logs pour scans QR',
    'Tests end-to-end pour flux réservation + scanning',
    'Améliorations UI: animations, toasts, accessibilité'
])

# Conclusion
add_bullet_slide('Conclusion', [
    'Application complète pour gestion d\'un petit cinéma',
    'Approche prudente pour migrations et sécurité des tokens',
    'Extensible: intégration paiement, notifications, analytics'
])

# Contact / références
add_bullet_slide('Contact & Références', [
    'Repo: (local project)',
    'Django docs, python-pptx, qrcode',
    'Questions ?'
])

# Save
prs.save(OUTFILE)
print(f'Présentation générée: {OUTFILE}')
